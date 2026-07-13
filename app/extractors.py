import csv
import io
import os
import subprocess
import tempfile
import zipfile
from pathlib import Path

import pytesseract
from PIL import Image
from docx import Document
from openpyxl import load_workbook
from pdf2image import convert_from_bytes
from pptx import Presentation
from pypdf import PdfReader

# Formats read directly, no conversion needed.
DIRECT = {'.docx', '.xlsx', '.pptx', '.pdf', '.txt', '.csv'}

# Legacy / other-suite formats: converted via headless LibreOffice into a
# DIRECT-readable equivalent first, then read with the same reader.
# ext -> target format LibreOffice should convert to
LEGACY_TARGET = {
    '.doc': 'docx',
    '.rtf': 'docx',
    '.odt': 'docx',
    '.xls': 'xlsx',
    '.ods': 'xlsx',
    '.ppt': 'pptx',
    '.odp': 'pptx',
}
LEGACY = set(LEGACY_TARGET)

ARCHIVE = {'.zip'}

IMAGES = {
    x.strip().lower()
    for x in os.getenv(
        'OCR_IMAGE_EXTENSIONS', '.jpg,.jpeg,.png,.tif,.tiff,.bmp,.webp'
    ).split(',')
    if x.strip()
}

CONTENT_EXTENSIONS = DIRECT | LEGACY | IMAGES | ARCHIVE

OCR = os.getenv('ENABLE_OCR', 'true').lower() in {'1', 'true', 'yes', 'on'}
LANG = os.getenv('OCR_LANG', 'tha+eng')
MAXP = max(1, int(os.getenv('OCR_MAX_PAGES', '20')))
ZIP_MAX_ENTRY_MB = int(os.getenv('ZIP_MAX_ENTRY_MB', '20'))
ZIP_MAX_ENTRIES = int(os.getenv('ZIP_MAX_ENTRIES', '200'))


def docx(data):
    d = Document(io.BytesIO(data))
    parts = [x.text for x in d.paragraphs if x.text.strip()]

    for t in d.tables:
        for r in t.rows:
            parts.append(' | '.join(c.text for c in r.cells))

    # Headers/footers repeat on every page and often carry case titles,
    # reference numbers, etc. that don't appear anywhere in the body text.
    for section in d.sections:
        for part in (section.header, section.footer):
            for p in part.paragraphs:
                if p.text.strip():
                    parts.append(p.text)
            for t in part.tables:
                for r in t.rows:
                    row_text = ' | '.join(c.text for c in r.cells)
                    if row_text.strip():
                        parts.append(row_text)

    return '\n'.join(parts)


def xlsx(data):
    w = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for s in w.worksheets:
        parts.append(f'[ชีต: {s.title}]')
        for r in s.iter_rows(values_only=True):
            v = [str(x) for x in r if x is not None]
            if v:
                parts.append(' | '.join(v))
    return '\n'.join(parts)


def pptx_text(data):
    pres = Presentation(io.BytesIO(data))
    parts = []
    for idx, slide in enumerate(pres.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    text = ''.join(run.text for run in p.runs)
                    if text.strip():
                        slide_parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = ' | '.join(cell.text for cell in row.cells)
                    if row_text.strip():
                        slide_parts.append(row_text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text
            if note.strip():
                slide_parts.append(f'[บันทึกย่อ] {note}')
        if slide_parts:
            parts.append(f'[สไลด์ {idx}]')
            parts.extend(slide_parts)
    return '\n'.join(parts)


def _convert_with_libreoffice(data, src_ext, target_fmt):
    with tempfile.TemporaryDirectory() as td:
        td = Path(td)
        src = td / f'source{src_ext}'
        src.write_bytes(data)
        r = subprocess.run(
            ['libreoffice', '--headless', '--convert-to', target_fmt, '--outdir', str(td), str(src)],
            capture_output=True,
            timeout=120,
        )
        if r.returncode != 0:
            raise RuntimeError(r.stderr.decode(errors='replace')[:500] or 'LibreOffice conversion failed')
        matches = list(td.glob(f'*.{target_fmt}'))
        if not matches:
            raise RuntimeError('Converted file was not created')
        return matches[0].read_bytes()


def legacy(data, ext):
    target = LEGACY_TARGET[ext]
    converted = _convert_with_libreoffice(data, ext, target)
    if target == 'docx':
        return docx(converted)
    if target == 'xlsx':
        return xlsx(converted)
    if target == 'pptx':
        return pptx_text(converted)
    raise ValueError(f'Unhandled LibreOffice conversion target: {target}')


def zip_archive(data):
    """Recurse into a .zip and extract text from every supported entry.

    Nested .zip entries are skipped (no recursive zip-bomb risk), and any
    single entry over ZIP_MAX_ENTRY_MB is skipped with a note instead of
    failing the whole archive.
    """
    parts = []
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        entries = [i for i in zf.infolist() if not i.is_dir()][:ZIP_MAX_ENTRIES]
        for info in entries:
            name = info.filename
            ext = Path(name).suffix.lower()
            if ext not in CONTENT_EXTENSIONS or ext in ARCHIVE:
                continue
            if info.file_size > ZIP_MAX_ENTRY_MB * 1024 * 1024:
                parts.append(f'[ไฟล์ใน ZIP: {name}] ข้ามเพราะไฟล์ใหญ่เกิน {ZIP_MAX_ENTRY_MB} MB')
                continue
            try:
                inner_text = extract_text(zf.read(info), ext).strip()
                if inner_text:
                    parts.append(f'[ไฟล์ใน ZIP: {name}]\n{inner_text}')
            except Exception as e:
                parts.append(f'[ไฟล์ใน ZIP: {name}] อ่านไม่ได้: {str(e)[:200]}')
    return '\n'.join(parts)


def extract_text(data, ext):
    ext = ext.lower()
    if ext == '.docx':
        return docx(data)
    if ext == '.xlsx':
        return xlsx(data)
    if ext == '.pptx':
        return pptx_text(data)
    if ext == '.pdf':
        txt = '\n'.join((p.extract_text() or '') for p in PdfReader(io.BytesIO(data)).pages).strip()
        if txt:
            return txt
        if not OCR:
            return ''
        return '\n'.join(
            pytesseract.image_to_string(im, lang=LANG)
            for im in convert_from_bytes(data, dpi=180, first_page=1, last_page=MAXP)
        )
    if ext == '.txt':
        return data.decode('utf-8', errors='replace')
    if ext == '.csv':
        return '\n'.join(' | '.join(r) for r in csv.reader(io.StringIO(data.decode('utf-8-sig', errors='replace'))))
    if ext in LEGACY:
        return legacy(data, ext)
    if ext in ARCHIVE:
        return zip_archive(data)
    if ext in IMAGES:
        return pytesseract.image_to_string(Image.open(io.BytesIO(data)), lang=LANG) if OCR else ''
    raise ValueError(f'Unsupported content extension: {ext}')
